import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(r"e:\Study\Competition\ICT\sign_language\ppt")
MD_PATH = ROOT / "内容.md"
TEMPLATE_PATH = ROOT / "【赛题1-模板2】创新赛-答辩PPT-学校-队名-队长姓名-手机号（中国总决赛阶段）-详细版.pptx"
OUTPUT_PATH = ROOT / "【赛题1-模板2】创新赛-答辩PPT-手语智瞳-张佳宇-16635829953（自动填充）.pptx"


SECTION_RE = re.compile(r"^##\s+第(\d+)页\s+(.+?)$", re.MULTILINE)


def parse_sections(md_text: str):
    matches = list(SECTION_RE.finditer(md_text))
    sections = []

    for i, m in enumerate(matches):
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(md_text)
        block = md_text[start:end].strip()
        page_no = int(m.group(1))
        page_title = m.group(2).strip()

        title = extract_single(block, r"^###\s+页面标题\s*\n(.+?)$", default=page_title)
        body = extract_bullets_or_paragraph(block, "页面正文（可直接放PPT）")
        script = extract_paragraph(block, "讲解词（可直接念）")

        sections.append({
            "page_no": page_no,
            "title": title,
            "body": body,
            "script": script,
        })

    return sections


def extract_single(text: str, pattern: str, default: str = ""):
    m = re.search(pattern, text, flags=re.MULTILINE)
    if not m:
        return default
    return m.group(1).strip()


def extract_block_by_header(text: str, header: str):
    # Match from target header to next level-3 header or end of section.
    pat = rf"^###\s+{re.escape(header)}\s*\n(.*?)(?=\n###\s+|\Z)"
    m = re.search(pat, text, flags=re.MULTILINE | re.DOTALL)
    if not m:
        return ""
    return m.group(1).strip()


def extract_bullets_or_paragraph(text: str, header: str):
    block = extract_block_by_header(text, header)
    if not block:
        return ""
    lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
    cleaned = []
    for ln in lines:
        if ln.startswith("-"):
            cleaned.append("• " + ln.lstrip("-").strip())
        else:
            cleaned.append(ln)
    return "\n".join(cleaned)


def extract_paragraph(text: str, header: str):
    block = extract_block_by_header(text, header)
    if not block:
        return ""
    return " ".join([ln.strip() for ln in block.splitlines() if ln.strip()])


def set_text_frame(tf, text: str, font_size=20, bold=False):
    tf.clear()
    paragraphs = [p for p in text.split("\n") if p.strip()]
    if not paragraphs:
        return

    first = tf.paragraphs[0]
    first.text = paragraphs[0]
    if first.runs:
        first.runs[0].font.size = Pt(font_size)
        first.runs[0].font.bold = bold

    for ptext in paragraphs[1:]:
        p = tf.add_paragraph()
        p.text = ptext
        if p.runs:
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.bold = bold


def ensure_notes(slide, note_text: str):
    if not note_text:
        return
    notes = slide.notes_slide
    notes.notes_text_frame.clear()
    notes.notes_text_frame.text = note_text


def fill_slide(slide, title: str, body: str, note: str):
    title_set = False
    body_set = False

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        if getattr(shape, "is_placeholder", False):
            ptype = shape.placeholder_format.type
            # title / center title
            if ptype in (1, 3) and not title_set:
                set_text_frame(shape.text_frame, title, font_size=36, bold=True)
                title_set = True
                continue
            # body / content
            if ptype in (2, 7) and not body_set:
                set_text_frame(shape.text_frame, body, font_size=20, bold=False)
                body_set = True
                continue

    if not title_set:
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.9))
        set_text_frame(title_box.text_frame, title, font_size=36, bold=True)

    if not body_set:
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.1))
        set_text_frame(body_box.text_frame, body, font_size=20, bold=False)

    ensure_notes(slide, note)


def main():
    md_text = MD_PATH.read_text(encoding="utf-8")
    sections = parse_sections(md_text)

    prs = Presentation(str(TEMPLATE_PATH))

    # Ensure enough slides
    while len(prs.slides) < len(sections):
        # Prefer a blank-like layout when available, otherwise use last layout.
        if len(prs.slide_layouts) >= 7:
            layout = prs.slide_layouts[6]
        else:
            layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
        prs.slides.add_slide(layout)

    for i, sec in enumerate(sections):
        slide = prs.slides[i]
        fill_slide(slide, sec["title"], sec["body"], sec["script"])

    prs.save(str(OUTPUT_PATH))

    print(f"SLIDES_FILLED={len(sections)}")
    print(f"TEMPLATE_SLIDES={len(prs.slides)}")
    print(f"OUTPUT={OUTPUT_PATH}")


if __name__ == "__main__":
    main()
